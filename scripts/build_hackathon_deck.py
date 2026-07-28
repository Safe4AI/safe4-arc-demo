"""Build the editable Safe4 Encode x Arc hackathon deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "artifacts" / "Safe4_Encode_Arc_Deck.pptx"

W = 13.333
H = 7.5

BG = "090B0A"
PANEL = "111512"
PANEL_2 = "171C18"
YELLOW = "E8FF3F"
WHITE = "F4F6F0"
MUTED = "99A39A"
LINE = "313A32"
GREEN = "70F6A5"
RED = "FF7A68"
CYAN = "7FE8FF"

FONT = "Arial"
MONO = "Consolas"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 20,
    fill: str = WHITE,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    uppercase: bool = False,
    tracking: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)
    if tracking is not None:
        run.font._element.set("spc", str(int(tracking * 1000)))
    return box


def add_rich_lines(
    slide,
    lines: list[tuple[str, str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    spacing: float = 10,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, (text, fill, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color(fill)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = PANEL,
    stroke: str = LINE,
    width: float = 1,
    radius: bool = False,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(stroke)
    shape.line.width = Pt(width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    stroke: str = LINE,
    width: float = 1,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color(stroke)
    line.line.width = Pt(width)
    return line


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = YELLOW,
    text_fill: str = BG,
):
    add_rect(slide, x, y, w, 0.34, fill=fill, stroke=fill, radius=True)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.25,
        size=9,
        fill=text_fill,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        uppercase=True,
    )


def add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str,
    title: str,
    body: str,
    accent: str = YELLOW,
):
    add_rect(slide, x, y, w, h, fill=PANEL, stroke=LINE)
    add_text(slide, label, x + 0.24, y + 0.2, w - 0.48, 0.24, size=9, fill=accent, bold=True, uppercase=True)
    add_text(slide, title, x + 0.24, y + 0.58, w - 0.48, 0.6, size=20, fill=WHITE, bold=True, uppercase=True)
    add_text(slide, body, x + 0.24, y + 1.25, w - 0.48, h - 1.48, size=12, fill=MUTED)


def add_base(slide, number: int, section: str, source: str = ""):
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = color(BG)

    # Quiet technical grid: editable vectors, not a flattened background.
    for x in (0.35, 4.55, 8.75, 12.95):
        add_line(slide, x, 0.15, x, 7.15, stroke="151A16", width=0.5)
    for y in (0.85, 6.9):
        add_line(slide, 0.35, y, 12.95, y, stroke=LINE, width=0.8)

    add_text(slide, "SAFE4", 0.42, 0.26, 1.2, 0.3, size=11, fill=YELLOW, bold=True, uppercase=True)
    add_text(slide, section, 1.7, 0.26, 4.6, 0.3, size=9, fill=MUTED, uppercase=True)
    add_text(slide, f"{number:02}", 12.25, 0.22, 0.62, 0.34, size=12, fill=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    if source:
        add_text(slide, source, 0.42, 7.05, 11.9, 0.18, size=6.5, fill="697169")


def add_title(slide, eyebrow: str, title: str, subtitle: str | None = None):
    add_text(slide, eyebrow, 0.42, 1.06, 3.4, 0.28, size=10, fill=YELLOW, bold=True, uppercase=True)
    add_text(slide, title, 0.42, 1.38, 12.0, 0.92, size=31, fill=WHITE, bold=True, uppercase=True)
    if subtitle:
        add_text(slide, subtitle, 0.44, 2.3, 11.7, 0.55, size=14, fill=MUTED)


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, accent: str = YELLOW):
    add_rect(slide, x, y, w, 1.12, fill=PANEL, stroke=LINE)
    add_text(slide, value, x + 0.22, y + 0.15, w - 0.44, 0.44, size=24, fill=accent, bold=True)
    add_text(slide, label, x + 0.22, y + 0.68, w - 0.44, 0.23, size=9, fill=MUTED, uppercase=True)


def build_deck() -> Presentation:
    deck = Presentation()
    deck.slide_width = Inches(W)
    deck.slide_height = Inches(H)
    blank = deck.slide_layouts[6]

    # 01 — Cover
    slide = deck.slides.add_slide(blank)
    add_base(slide, 1, "Encode x Arc · Agentic Economy")
    add_pill(slide, "Payment security / Arc Testnet", 0.44, 1.17, 2.65)
    add_text(slide, "THE PAYMENT\nFIREWALL FOR\nAI AGENTS", 0.4, 1.82, 8.7, 2.8, size=45, fill=WHITE, bold=True, uppercase=True)
    add_text(
        slide,
        "The security layer between an agent deciding to spend—and the money moving.",
        0.46,
        4.96,
        7.3,
        0.72,
        size=18,
        fill=MUTED,
    )
    add_rect(slide, 9.3, 1.3, 3.1, 4.65, fill=PANEL, stroke=YELLOW, width=1.3)
    add_text(slide, "SAFE", 9.68, 1.77, 2.3, 0.55, size=26, fill=YELLOW, bold=True, uppercase=True)
    add_text(slide, "04", 9.65, 2.38, 2.3, 1.45, size=75, fill=WHITE, bold=True)
    add_line(slide, 9.68, 4.06, 11.98, 4.06, stroke=YELLOW, width=1)
    add_text(slide, "POLICY\nPROOF\nPURPOSE", 9.68, 4.34, 2.3, 1.18, size=14, fill=MUTED, bold=True)

    # 02 — Problem
    slide = deck.slides.add_slide(blank)
    add_base(slide, 2, "Problem")
    add_title(slide, "The security gap", "The missing control point", "Documented wallet controls do not evaluate the task-to-purchase match demonstrated here.")
    add_panel(slide, 0.44, 3.0, 3.75, 2.65, label="01 / Decision", title="Agent chooses", body="A model interprets a task, selects a service, and proposes a payment.", accent=CYAN)
    add_panel(slide, 4.78, 3.0, 3.75, 2.65, label="02 / Gap", title="Context disappears", body="Amount and address survive. The assigned task, purchase purpose, and autonomy scope often do not.", accent=RED)
    add_panel(slide, 9.12, 3.0, 3.75, 2.65, label="03 / Execution", title="Money moves", body="A valid wallet action can still be the wrong action for the job.", accent=YELLOW)
    add_line(slide, 4.2, 4.33, 4.76, 4.33, stroke=CYAN, width=2)
    add_line(slide, 8.54, 4.33, 9.1, 4.33, stroke=YELLOW, width=2)
    add_text(slide, "SAFE4 INSERTS AN ENFORCEABLE CHECKPOINT HERE", 3.98, 5.96, 5.5, 0.34, size=11, fill=YELLOW, bold=True, align=PP_ALIGN.CENTER, uppercase=True)

    # 03 — Why now
    slide = deck.slides.add_slide(blank)
    add_base(
        slide,
        3,
        "Market timing",
        "Sources: x402.org · cloud.google.com/blog/products/ai-machine-learning/announcing-agents-to-payments-ap2-protocol · developers.circle.com/agent-stack",
    )
    add_title(slide, "Market shift", "Why now", "Agent-native payment rails are arriving faster than agent-native security controls.")
    add_metric(slide, "x402", "HTTP-native payment requirements", 0.44, 3.1, 3.72, CYAN)
    add_metric(slide, "AP2", "Mandates and agent payment intent", 4.8, 3.1, 3.72, YELLOW)
    add_metric(slide, "USDC", "Programmable settlement asset", 9.16, 3.1, 3.72, GREEN)
    add_rect(slide, 0.44, 4.62, 12.44, 1.25, fill=PANEL_2, stroke=YELLOW)
    add_text(slide, "New rails make autonomous commerce possible.", 0.73, 4.93, 5.35, 0.36, size=18, fill=WHITE, bold=True)
    add_text(slide, "Safe4 makes the payment decision inspectable before execution.", 6.2, 4.93, 6.32, 0.36, size=18, fill=YELLOW, bold=True)

    # 04 — Evaluation
    slide = deck.slides.add_slide(blank)
    add_base(slide, 4, "Policy engine")
    add_title(slide, "Decision surface", "What Safe4 evaluates", "One authorization decision, multiple independent checks, one durable reason trail.")
    checks = [
        ("01", "TASK ↔ PURCHASE", "Does the purpose match the submitted task context?", YELLOW),
        ("02", "AMOUNT + BUDGET", "Is it within transaction, daily, and velocity limits?", GREEN),
        ("03", "COUNTERPARTY", "Is the vendor or address blocked or otherwise risky?", CYAN),
        ("04", "AUTONOMY SCOPE", "Is the agent allowed to make this class of decision?", WHITE),
        ("05", "PAYMENT PROOF", "Does the x402 receipt match amount, currency, and request?", YELLOW),
        ("06", "AUDITABILITY", "Can the decision and evidence be reconstructed later?", GREEN),
    ]
    positions = [(0.44, 2.85), (4.56, 2.85), (8.68, 2.85), (0.44, 4.6), (4.56, 4.6), (8.68, 4.6)]
    for (num, heading, body, accent), (x, y) in zip(checks, positions):
        add_rect(slide, x, y, 3.72, 1.42, fill=PANEL, stroke=LINE)
        add_text(slide, num, x + 0.2, y + 0.18, 0.38, 0.25, size=9, fill=accent, bold=True)
        add_text(slide, heading, x + 0.62, y + 0.16, 2.82, 0.3, size=13, fill=WHITE, bold=True)
        add_text(slide, body, x + 0.2, y + 0.62, 3.3, 0.52, size=10.5, fill=MUTED)

    # 05 — Architecture
    slide = deck.slides.add_slide(blank)
    add_base(slide, 5, "Architecture", "Observed Arc Testnet path: Safe4 ALLOW precedes Circle Agent Wallet execution.")
    add_title(slide, "Observed path", "Safe4 authorization + Circle Agent Wallet", "Fresh 0.01 testnet USDC settlement is RPC-verified through the ERC-4337 receipt.")
    nodes = [
        (0.42, "AGENT", "Task + budget"),
        (2.84, "x402 SERVICE", "Payment required"),
        (5.26, "SAFE4", "Policy + purpose"),
        (7.68, "CIRCLE", "Agent Wallet"),
        (10.1, "ARC", "USDC settlement"),
    ]
    for index, (x, label, sub) in enumerate(nodes):
        stroke = YELLOW if label == "SAFE4" else LINE
        add_rect(slide, x, 3.25, 1.92, 1.3, fill=PANEL, stroke=stroke, width=1.4 if label == "SAFE4" else 1)
        add_text(slide, f"0{index + 1}", x + 0.16, 3.43, 0.35, 0.22, size=9, fill=YELLOW, bold=True)
        add_text(slide, label, x + 0.16, 3.75, 1.6, 0.26, size=13, fill=WHITE, bold=True, uppercase=True)
        add_text(slide, sub, x + 0.16, 4.13, 1.6, 0.2, size=8.5, fill=MUTED)
        if index < len(nodes) - 1:
            add_line(slide, x + 1.94, 3.9, x + 2.38, 3.9, stroke=YELLOW, width=1.8)
    add_rect(slide, 5.0, 4.95, 2.45, 0.74, fill=PANEL_2, stroke=YELLOW)
    add_text(slide, "ALLOW → REPLAY VERIFY", 5.12, 5.12, 2.22, 0.22, size=9.5, fill=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 7.78, 4.95, 2.7, 0.74, fill=PANEL_2, stroke=RED)
    add_text(slide, "DENY → DEMO EXECUTOR SKIPPED", 7.84, 5.12, 2.58, 0.22, size=8.8, fill=RED, bold=True, align=PP_ALIGN.CENTER)

    # 06 — Evidence
    slide = deck.slides.add_slide(blank)
    add_base(
        slide,
        6,
        "Observed evidence",
        "Arcscan: testnet.arcscan.app/tx/0x648ef1…7752c · ERC-4337 verifier: scripts/verify_arc_settlement.py",
    )
    add_title(slide, "Observed output", "Real demo. Real reasons. Real chain evidence.")
    add_rect(slide, 0.44, 2.72, 5.92, 2.62, fill=PANEL, stroke=GREEN, width=1.4)
    add_pill(slide, "Allowed", 0.72, 2.98, 1.18, fill=GREEN)
    add_text(slide, "TASK_PURCHASE_MATCH", 0.72, 3.52, 5.1, 0.34, size=18, fill=WHITE, bold=True)
    add_text(slide, "Competitor-pricing research matches the assigned task and allowed category.", 0.72, 3.97, 5.05, 0.58, size=12.5, fill=MUTED)
    add_text(slide, "HISTORICAL REPLAY · NOT BROADCAST BY DEMO", 0.72, 4.78, 5.08, 0.25, size=9.2, fill=GREEN, bold=True)
    add_rect(slide, 6.94, 2.72, 5.92, 2.62, fill=PANEL, stroke=RED, width=1.4)
    add_pill(slide, "Denied", 7.22, 2.98, 1.18, fill=RED, text_fill=BG)
    add_text(slide, "PURCHASE_PURPOSE_MISMATCH", 7.22, 3.52, 5.1, 0.34, size=18, fill=WHITE, bold=True)
    add_text(slide, "Same amount, category, and counterparty. Gift-card purpose does not match the research task.", 7.22, 3.97, 5.0, 0.58, size=12.5, fill=MUTED)
    add_text(slide, "DEMO ORCHESTRATOR DID NOT INVOKE EXECUTOR", 7.22, 4.78, 5.08, 0.25, size=9.2, fill=RED, bold=True)
    add_rect(slide, 0.44, 5.75, 12.42, 0.7, fill=PANEL_2, stroke=LINE)
    add_text(slide, "RPC-VERIFIED HISTORICAL ARC TX", 0.67, 5.97, 2.35, 0.2, size=8.5, fill=YELLOW, bold=True)
    add_text(slide, "0x648ef14e4da7c6bfecce0017d19280ed51fb12635bea94712de926d9f967752c", 3.06, 5.91, 9.12, 0.28, size=9.2, fill=WHITE, font=MONO)

    # 07 — Arc / Circle rationale
    slide = deck.slides.add_slide(blank)
    add_base(
        slide,
        7,
        "Stack rationale",
        "Sources: docs.arc.io · developers.circle.com/agent-stack · developers.circle.com/agent-stack/agent-wallets/supported-blockchains",
    )
    add_title(slide, "Infrastructure choice", "Why Arc, USDC, and Circle")
    add_panel(slide, 0.44, 2.8, 3.75, 3.06, label="Arc", title="Payment-native chain", body="A purpose-built environment for programmable money. Safe4's exact verifier checks chain, token, calldata, receipt status, and Transfer event.", accent=YELLOW)
    add_panel(slide, 4.79, 2.8, 3.75, 3.06, label="USDC", title="Precise settlement", body="Six-decimal, dollar-denominated testnet evidence makes the demo amount and token contract unambiguous.", accent=GREEN)
    add_panel(slide, 9.14, 2.8, 3.75, 3.06, label="Circle Agent Stack", title="Agent-native execution", body="An authenticated Agent Wallet settled 0.01 testnet USDC after ALLOW. Safe4 verifies its ERC-4337 receipt.", accent=CYAN)

    # 08 — Differentiation
    slide = deck.slides.add_slide(blank)
    add_base(slide, 8, "Differentiation", "Circle control descriptions: developers.circle.com/agent-stack/agent-wallets")
    add_title(slide, "Complementary controls", "Circle enforces the floor. Safe4 asks whether the payment should happen.")
    add_rect(slide, 0.44, 2.8, 12.42, 2.95, fill=PANEL, stroke=LINE)
    add_text(slide, "CONTROL", 0.72, 3.03, 2.6, 0.25, size=9, fill=MUTED, bold=True)
    add_text(slide, "CIRCLE AGENT WALLET", 3.72, 3.03, 3.4, 0.25, size=9, fill=CYAN, bold=True)
    add_text(slide, "SAFE4", 8.02, 3.03, 3.8, 0.25, size=9, fill=YELLOW, bold=True)
    rows = [
        ("Spending limits", "✓ Native control", "✓ Additional policy context"),
        ("Address / compliance checks", "✓ Native guardrails", "✓ Counterparty + audit context"),
        ("Submitted task ↔ purchase", "Documented controls do not evaluate this match", "✓ Outcome-changing decision"),
    ]
    y = 3.55
    for label, circle_value, safe4_value in rows:
        add_line(slide, 0.72, y - 0.12, 12.5, y - 0.12, stroke=LINE)
        add_text(slide, label, 0.72, y, 2.65, 0.34, size=12, fill=WHITE, bold=True)
        add_text(slide, circle_value, 3.72, y, 3.48, 0.34, size=12, fill=MUTED)
        add_text(slide, safe4_value, 8.02, y, 4.1, 0.34, size=12, fill=YELLOW if "Outcome" in safe4_value else WHITE)
        y += 0.65
    add_text(slide, "DEMO MOMENT", 0.72, 6.04, 1.4, 0.23, size=9, fill=RED, bold=True)
    add_text(slide, "Same amount, category, counterparty. Denied against submitted task context. Trusted binding is roadmap work.", 2.16, 5.98, 10.1, 0.37, size=12.2, fill=WHITE, bold=True)

    # 09 — Team
    slide = deck.slides.add_slide(blank)
    add_base(slide, 9, "Team", "Team experience is a human-confirmed submission fact; names intentionally omitted until supplied.")
    add_title(slide, "Team", "Built at the intersection of cybersecurity and finance")
    add_text(slide, "35", 0.42, 2.67, 3.1, 1.42, size=78, fill=YELLOW, bold=True)
    add_text(slide, "COMBINED YEARS", 0.49, 4.18, 3.1, 0.34, size=12, fill=WHITE, bold=True, uppercase=True)
    add_text(slide, "IN REGULATED ENVIRONMENTS", 0.49, 4.58, 3.2, 0.32, size=10, fill=MUTED, uppercase=True)
    add_panel(slide, 4.35, 2.7, 3.8, 2.45, label="Security", title="Threat-first", body="Payment controls, attack surfaces, evidence, and operational resilience.", accent=YELLOW)
    add_panel(slide, 8.62, 2.7, 3.8, 2.45, label="Finance", title="Risk-aware", body="Money movement, governed operations, and decision accountability.", accent=GREEN)
    add_text(slide, "Founder names and biographies will be added only from confirmed source material.", 4.37, 5.62, 7.9, 0.42, size=11, fill=MUTED)

    # 10 — Progress / roadmap
    slide = deck.slides.add_slide(blank)
    add_base(slide, 10, "Progress + accelerator path", "Public repo: github.com/Safe4AI/safe4-arc-demo · evidence current at 28 Jul 2026")
    add_title(slide, "Roadmap", "From verified demo to agent-commerce service")
    add_panel(slide, 0.44, 2.75, 3.75, 2.85, label="Now / Verified", title="Hackathon proof", body="293-test regression gate. Fresh Circle Agent Wallet settlement on Arc. Unattended allow/deny demo. Public repository.", accent=GREEN)
    add_panel(slide, 4.79, 2.75, 3.75, 2.85, label="Next / Harden", title="Trusted context", body="Bind tasks to principals, remove legacy bypasses, harden the x402 verifier boundary, and complete security review.", accent=YELLOW)
    add_panel(slide, 9.14, 2.75, 3.75, 2.85, label="After / Distribute", title="Agent service", body="Package Safe4 as an x402 service, complete security review, and apply for Circle Agent Marketplace listing.", accent=CYAN)
    add_text(slide, "THE ACCELERATOR UNLOCKS", 0.47, 6.02, 2.3, 0.25, size=9, fill=YELLOW, bold=True)
    add_text(slide, "design partners · hardened integrations · marketplace readiness", 2.82, 5.95, 9.55, 0.34, size=15, fill=WHITE, bold=True)

    return deck


def main() -> int:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    deck = build_deck()
    deck.save(OUTPUT)
    print(f"DECK_OK path={OUTPUT} slides={len(deck.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
